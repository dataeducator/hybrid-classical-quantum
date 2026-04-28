"""Generate the final presentation following the MSDS 730 template structure.

Template sections:
  1. Title (project, team, course/date)
  2. Motivation & Problem Statement
  3. Related Work / Background
  4. Dataset & Preprocessing
  5. Methodology and Model Architecture
  6. Results
  7. Demo (optional)
  8. Discussion
  9. Conclusion & Future Work
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Meharry brand colors (consistent with paper)
MAROON = RGBColor(0x6C, 0x1A, 0x42)
TEAL = RGBColor(0x1A, 0x6C, 0x5E)
GOLD = RGBColor(0xC4, 0x9A, 0x2A)
BLUE = RGBColor(0x2C, 0x5F, 0x8A)
DARK = RGBColor(0x1F, 0x1F, 0x1F)
GRAY = RGBColor(0x55, 0x55, 0x55)

BLANK = prs.slide_layouts[6]


def add_slide(title=None, title_color=MAROON):
    slide = prs.slides.add_slide(BLANK)
    if title:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = title_color
    return slide


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=DARK, align=None, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(text, list):
        for i, line in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.italic = italic
            p.font.color.rgb = color
            if align:
                p.alignment = align
    else:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        if align:
            p.alignment = align
    return tb


def add_bullets(slide, items, left, top, width, height, size=16, color=DARK):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if level == 0:
            bullet = "• "
            p_size = size
            p_color = color
        elif level == 1:
            bullet = "    – "
            p_size = size - 2
            p_color = GRAY
        else:
            bullet = "        · "
            p_size = size - 3
            p_color = GRAY
        p.text = bullet + text
        p.font.size = Pt(p_size)
        p.font.color.rgb = p_color


def add_table(slide, data, left, top, width, height, header_color=MAROON,
              first_col_width_in=None):
    n_rows = len(data)
    n_cols = len(data[0])
    table = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                    Inches(width), Inches(height)).table
    if first_col_width_in:
        table.columns[0].width = Inches(first_col_width_in)
        remaining = (width - first_col_width_in) / max(1, n_cols - 1)
        for j in range(1, n_cols):
            table.columns[j].width = Inches(remaining)
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    if i == 0:
                        r.font.bold = True
                        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    else:
                        r.font.color.rgb = DARK
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    return table


# ============================================================
# Slide 1: Title
# ============================================================
slide = prs.slides.add_slide(BLANK)
# Background accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.6),
                               Inches(13.333), Inches(0.15))
bar.fill.solid(); bar.fill.fore_color.rgb = MAROON; bar.line.fill.background()

add_text(slide, "Hybrid Quantum-Classical Neural Networks for",
         0.5, 1.4, 12.3, 0.6, size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(slide, "Personalized Survival Prediction in Triple-Negative Breast Cancer",
         0.5, 1.95, 12.3, 0.6, size=28, bold=True, color=MAROON, align=PP_ALIGN.CENTER)

add_text(slide, "A Fairness-Aware, Survival-Analytic Modeling Framework",
         0.5, 3.0, 12.3, 0.5, size=20, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "Jaclyn Claiborne and Tenicka Norwood",
         0.5, 4.5, 12.3, 0.5, size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(slide, "MSDS 730: Deep Learning",
         0.5, 5.1, 12.3, 0.4, size=18, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Meharry Medical College | Spring 2026",
         0.5, 5.5, 12.3, 0.4, size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 2: Motivation & Problem Statement
# ============================================================
slide = add_slide("Motivation & Problem Statement")

add_text(slide, "Problem", 0.5, 1.4, 6, 0.5, size=20, bold=True, color=TEAL)
add_bullets(slide, [
    ("Triple-negative breast cancer (TNBC) is the most aggressive subtype.", 0),
    ("Predicting individual survival informs treatment but is hard with right-censored data.", 0),
    ("Existing models often perform worse on minority subpopulations.", 0),
], 0.5, 1.9, 6, 2.5)

add_text(slide, "Why it matters", 0.5, 4.5, 6, 0.5, size=20, bold=True, color=TEAL)
add_bullets(slide, [
    ("Health equity: Black women have 2x TNBC incidence and worse outcomes.", 0),
    ("Clinical decision support: better risk stratification informs care.", 0),
    ("Methodology: rigorous leakage auditing for clinical ML pipelines.", 0),
], 0.5, 5.0, 6, 2.0)

add_text(slide, "Stakeholders", 7, 1.4, 6, 0.5, size=20, bold=True, color=GOLD)
add_bullets(slide, [
    ("Oncologists making treatment-planning decisions.", 0),
    ("Patients receiving personalized risk estimates.", 0),
    ("Policy researchers studying disparities in clinical AI.", 0),
    ("ML practitioners building survival models on registry data.", 0),
], 7, 1.9, 6, 3)


# ============================================================
# Slide 3: Related Work
# ============================================================
slide = add_slide("Related Work & Background")

add_text(slide, "Existing approaches", 0.5, 1.4, 12, 0.5, size=20, bold=True, color=TEAL)
add_bullets(slide, [
    ("Cox proportional hazards (Cox 1972) — gold-standard for tabular survival.", 0),
    ("DeepSurv (Katzman et al. 2018) — neural Cox via partial likelihood.", 0),
    ("Variational quantum circuits for classification (Schuld 2021, Benedetti 2019).", 0),
    ("Data re-uploading universal classifiers (Pérez-Salinas et al. 2020).", 0),
    ("Algorithm fairness in clinical AI (Chen et al. 2023, Raza et al. 2024).", 0),
], 0.5, 1.9, 12, 3.0)

add_text(slide, "Gap we address", 0.5, 5.0, 12, 0.5, size=20, bold=True, color=TEAL)
add_bullets(slide, [
    ("First TNBC survival study using hybrid quantum-classical with Cox loss.", 0),
    ("Automated leakage auditing as a methodology contribution.", 0),
    ("Subgroup C-index fairness for survival models (vs naive binary fairness).", 0),
    ("Mathematically guaranteed quantum residual learning framework (≥ Cox PH).", 0),
], 0.5, 5.5, 12, 2)


# ============================================================
# Slide 4: Dataset
# ============================================================
slide = add_slide("Dataset: SEER TNBC Registry")

add_text(slide, "Source", 0.5, 1.4, 6, 0.5, size=20, bold=True, color=TEAL)
add_bullets(slide, [
    ("SEER Surveillance, Epidemiology, and End Results", 0),
    ("https://seer.cancer.gov/data/", 1),
    ("Filtered to TNBC: ER−, PR−, HER2−", 0),
], 0.5, 1.9, 6, 1.6)

add_text(slide, "Cohort sizes", 0.5, 3.7, 6, 0.5, size=20, bold=True, color=TEAL)
add_bullets(slide, [
    ("Raw export: 31,683 patients", 0),
    ("After right-censoring fix: 20,110 patients", 0),
    ("• 4,022 events (deaths within 60 mo)", 1),
    ("• 16,088 confirmed survivors at 60 mo", 1),
    ("Train / test: 80/20 stratified, seed=42", 0),
], 0.5, 4.2, 6, 3)

add_text(slide, "Features (27 total)", 7, 1.4, 6, 0.5, size=20, bold=True, color=GOLD)
add_bullets(slide, [
    ("7 quantum features: stage, age, tumor size, income, treatment timing, etc.", 0),
    ("20 classical features: race (6 cat), marital, subtype, surgery, chemo, etc.", 0),
    ("Modality: structured tabular clinical/demographic", 0),
], 7, 1.9, 6, 2)

add_text(slide, "Class distribution (binary @ 60mo)", 7, 4.5, 6, 0.5, size=20, bold=True, color=GOLD)
add_bullets(slide, [
    ("76% positive (alive at 60+ mo) | 24% negative (deceased)", 0),
    ("Severely imbalanced → class-weighted loss for neural models", 0),
], 7, 5, 6, 2)


# ============================================================
# Slide 5: Preprocessing & Leakage Audit
# ============================================================
slide = add_slide("Preprocessing & Data Leakage Audit")

add_text(slide, "Cleaning & normalization", 0.5, 1.3, 12, 0.5, size=20, bold=True, color=TEAL)
add_bullets(slide, [
    ("Tumor size: backfilled from CS_tumor_size column, then median-filled.", 0),
    ("Race encoding: 6 SEER categories preserved (incl Non-Hispanic Unknown).", 0),
    ("Quantum features: MinMax scaled to [0, π] for RY rotation encoding.", 0),
    ("Classical features: StandardScaler (μ=0, σ=1).", 0),
    ("Right-censoring: drop alive patients with <60 mo follow-up (label noise).", 0),
], 0.5, 1.8, 12, 2.5)

add_text(slide, "Leakage audit (3 automated checks)", 0.5, 4.5, 12, 0.5,
         size=20, bold=True, color=MAROON)
add_bullets(slide, [
    ("Per-feature AUC > 0.95: flagged Year_of_follow_up_recode (AUC=0.98). EXCLUDED.", 0),
    ("Feature-target correlation > 0.90: PASS.", 0),
    ("Train/test KS distribution shift: PASS.", 0),
    ("→ Leak removal alone dropped pilot AUC from 0.92 → ~0.74 (honest).", 0),
], 0.5, 5, 12, 2)


# ============================================================
# Slide 6: Methodology - Architecture
# ============================================================
slide = add_slide("Methodology: Hybrid Architecture")

add_text(slide, "Three model generations + residual framework", 0.5, 1.3, 12, 0.5,
         size=20, bold=True, color=TEAL)
add_bullets(slide, [
    ("v1 (baseline): 7-qubit VQC, single-layer RY, ring CNOT, single Z-measurement → 1D scalar.", 0),
    ("v2: 3-layer VQC + data re-uploading + RY/RZ + shifted-ring CNOT + all-qubit Z → 7D embedding.", 0),
    ("v3: v2 + trainable input scaling + ZZ correlator measurements → 14D embedding.", 0),
    ("v4: v3 + classical pre-training + learnable output scale/bias + smaller init.", 0),
    ("Cox + Quantum Residual: fixed Cox PH offset + small quantum correction (the winner).", 0),
], 0.5, 1.8, 12, 3)

# Add architecture image if available
arch_path = "figures/hybrid_architecture.png"
if not os.path.exists(arch_path):
    # Try the PDF rendered version
    arch_path = "figures/tikz_architecture.png"
if os.path.exists(arch_path):
    try:
        slide.shapes.add_picture(arch_path, Inches(8), Inches(4.5), width=Inches(5))
    except Exception:
        pass

add_text(slide, "Residual learning equation", 0.5, 5, 7, 0.5,
         size=18, bold=True, color=MAROON)
add_text(slide, "h_total(x) = h_Cox(x) + s · q_θ(x)",
         0.5, 5.5, 7, 0.5, size=22, bold=True, color=DARK)
add_text(slide, "At init: q_θ ≈ 0 → model = Cox PH. Training only improves on Cox.",
         0.5, 6.0, 7, 1.0, size=14, italic=True, color=GRAY)


# ============================================================
# Slide 7: Methodology - Training
# ============================================================
slide = add_slide("Methodology: Training Details")

table_data = [
    ["Component", "Choice"],
    ["Loss (survival)", "Cox partial likelihood (negative log)"],
    ["Loss (binary)", "BCEWithLogitsLoss with class pos_weight"],
    ["Optimizer", "Adam, lr=0.001 (lr=0.005 for v4 pretrain)"],
    ["LR schedule", "CosineAnnealingLR (1e-3 → 1e-5) for v2-v4 + Residual"],
    ["Epochs", "50 (v1) | 75 (v2-v4) | 15 (Residual) | 20 (full-data)"],
    ["Batch size", "Full-batch on 2k subsample (v1-v3) | 256 mini-batch on full data"],
    ["Quantum simulator", "PennyLane lightning.qubit (CPU)"],
    ["Hardware", "NVIDIA RTX 5080, PyTorch 2.10.0, CUDA 12.8"],
    ["Seed", "42 across Python/NumPy/PyTorch/CUDA"],
]
add_table(slide, table_data, 1.5, 1.4, 10, 5.5,
          header_color=MAROON, first_col_width_in=3.5)


# ============================================================
# Slide 8: Results - Survival Ablation
# ============================================================
slide = add_slide("Results: Survival Ablation (C-index)")

surv_table = [
    ["Model", "Test C-index", "Subgroup Gap", "Time (s)", "Train N"],
    ["Cox PH (lifelines)", "0.7326", "0.0842", "0.2", "16,088"],
    ["HybridSurvivalQ_v1 (subsample)", "0.5869", "—", "223", "2,000"],
    ["HybridSurvivalQ_v2 (subsample)", "0.5759", "—", "575", "2,000"],
    ["HybridSurvivalQ_v3 (subsample)", "0.5872", "—", "615", "2,000"],
    ["HybridSurvivalQ_v1 (full data)", "0.6897", "—", "651", "16,088"],
    ["HybridSurvivalQ_v2 (full data)", "0.7192", "—", "1,280", "16,088"],
    ["HybridSurvivalQ_v3 (full data)", "0.6674", "0.1614", "1,038", "16,088"],
    ["HybridSurvivalQ_v4 (full data)", "0.7249", "0.0371", "1,321", "16,088"],
    ["Cox + Quantum Residual ★", "0.7364", "0.0754", "1,237", "16,088"],
]
add_table(slide, surv_table, 0.5, 1.4, 12.3, 5,
          header_color=MAROON, first_col_width_in=4.5)

add_text(slide, "★ Only model that beats classical baseline (+0.0038 over Cox PH).",
         0.5, 6.6, 12, 0.5, size=14, italic=True, bold=True, color=MAROON)


# ============================================================
# Slide 9: Results - Binary Ablation
# ============================================================
slide = add_slide("Results: Binary Ablation (AUC at 60 months)")

bin_table = [
    ["Model", "AUC", "Precision", "Recall", "F1", "Train N"],
    ["LightGBM (full, balanced)", "0.7462", "0.8715", "0.7173", "0.7869", "15,673"],
    ["XGBoost (full, balanced)", "0.7434", "0.8750", "0.7009", "0.7784", "15,673"],
    ["Cox + Residual @ 60mo (cross-fwk)", "0.7426", "0.8594", "0.7255", "0.7868", "15,673"],
    ["Cox PH @ 60mo (cross-fwk)", "0.7402", "0.8668", "0.7072", "0.7789", "15,673"],
    ["HybridRealQ_v4 (full, output scale)", "0.6863", "0.8535", "0.6088", "0.7107", "15,673"],
    ["HybridRealQ_v3 (full data)", "0.6828", "0.8709", "0.5090", "0.6425", "15,673"],
    ["Ensemble (v3-full + MLP)", "0.6770", "0.8728", "0.4740", "0.6143", "15,673"],
    ["HybridRealQ_v2 (subsample)", "0.5923", "0.8704", "0.3588", "0.5081", "2,000"],
    ["Classical MLP (full, weighted)", "0.5792", "0.8120", "0.5434", "0.6511", "15,673"],
    ["HybridRealQ_v3 (subsample)", "0.5718", "0.8137", "0.8779", "0.8446", "2,000"],
    ["HybridRealQ_v1 (subsample)", "0.5649", "0.8228", "0.6972", "0.7548", "2,000"],
]
add_table(slide, bin_table, 0.5, 1.4, 12.3, 5.2,
          header_color=BLUE, first_col_width_in=4.5)

add_text(slide, "Key finding: Cox+Residual evaluated at 60mo (AUC 0.7426) is competitive with LightGBM (0.7462). Cross-framework consistency.",
         0.5, 6.7, 12, 0.5, size=13, italic=True, color=GRAY)


# ============================================================
# Slide 10: Results - Fairness
# ============================================================
slide = add_slide("Results: Subgroup Fairness Audit")

fair_table = [
    ["Subgroup", "N", "Events", "Cox PH C-idx", "Cox + Residual C-idx"],
    ["Non-Hispanic White", "2,218", "598", "0.7267", "0.7300"],
    ["Non-Hispanic Black", "599", "228", "0.7195", "0.7176"],
    ["Hispanic", "716", "218", "0.7119", "0.7263"],
    ["Non-Hispanic Asian/PI", "435", "106", "0.7521", "0.7486"],
    ["Non-Hispanic AIAN*", "36", "11", "0.7962", "0.7930"],
    ["Non-Hispanic Unknown†", "18", "1", "—", "—"],
    ["Subgroup gap (max-min)", "", "", "0.0842", "0.0754"],
]
add_table(slide, fair_table, 0.5, 1.4, 12, 4.5,
          header_color=GOLD, first_col_width_in=4)

add_text(slide, "* Small sample (11 events) → high variance; † N<20 events<5 → audit excluded",
         0.5, 6.0, 12, 0.4, size=12, italic=True, color=GRAY)
add_text(slide, "Survival framing produces more uniform subgroup performance than binary framing.",
         0.5, 6.4, 12, 0.4, size=14, bold=True, italic=True, color=DARK)
add_text(slide, "All major subgroups achieve C-index ≥ 0.71. Residual model tightens the gap (0.0842 → 0.0754).",
         0.5, 6.8, 12, 0.4, size=14, color=GRAY)


# ============================================================
# Slide 11: Results - Visualization
# ============================================================
slide = add_slide("Results: Training Dynamics")

# Try to embed the survival training curve
img_path = "figures/survival_training_curve.png"
if os.path.exists(img_path):
    try:
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.4), width=Inches(8))
    except Exception:
        pass

add_text(slide, "Two-phase v4 training",
         9, 1.5, 4, 0.5, size=18, bold=True, color=TEAL)
add_bullets(slide, [
    ("Phase 1: classical pretrain plateaus at 0.56", 0),
    ("Phase 2: quantum branch enabled → jumps to 0.70 in 4 epochs", 0),
    ("Final: 0.7249 just below Cox PH 0.7326", 0),
], 9, 2.0, 4.2, 2.5, size=13)

add_text(slide, "Residual learning",
         9, 4.7, 4, 0.5, size=18, bold=True, color=MAROON)
add_bullets(slide, [
    ("Starts at Cox PH C-idx 0.7326 by construction", 0),
    ("Climbs steadily to 0.7364 over 9 epochs", 0),
    ("Quantum correction grows σ: 0.005 → 0.205", 0),
], 9, 5.2, 4.2, 2, size=13)


# ============================================================
# Slide 12: Discussion - What Worked
# ============================================================
slide = add_slide("Discussion: What Worked & Didn't")

add_text(slide, "What worked", 0.5, 1.3, 6, 0.5, size=20, bold=True, color=TEAL)
add_bullets(slide, [
    ("Automated leakage auditing caught a critical leak.", 0),
    ("Survival framing avoided right-censoring label noise.", 0),
    ("Output scale layer in v4: AUC ~0.58 → 0.72.", 0),
    ("Mini-batch on full data: subsample 0.59 → 0.69 (v1) and 0.58 → 0.72 (v2).", 0),
    ("Quantum residual learning beats Cox PH on both metrics.", 0),
    ("Subgroup C-index gap tightened by residual learning.", 0),
], 0.5, 1.8, 6, 4.5, size=14)

add_text(slide, "What didn't work", 7, 1.3, 6, 0.5, size=20, bold=True, color=MAROON)
add_bullets(slide, [
    ("Subsample-only neural models (v1-v3) couldn't escape AUC 0.55-0.60.", 0),
    ("Naive ensembling (v3-full + MLP): worse than v3-full alone.", 0),
    ("Classical MLP on imbalanced data: AUC 0.58 even with class weighting.", 0),
    ("Quantum-only architecture: cannot exploit classical features.", 0),
    ("Per-sample quantum simulation: limits epoch count and grid search.", 0),
], 7, 1.8, 6, 4.5, size=14)


# ============================================================
# Slide 13: Discussion - Challenges & Limitations
# ============================================================
slide = add_slide("Key Challenges & Limitations")

add_text(slide, "Key challenges encountered", 0.5, 1.3, 12, 0.5,
         size=20, bold=True, color=TEAL)
add_bullets(slide, [
    ("Discovering Year_of_follow_up_recode leakage forced complete re-analysis.", 0),
    ("Switching from binary to survival framing required new loss + new evaluation.", 0),
    ("Cox loss has flat gradients without output rescaling → v4 output scale fix.", 0),
    ("Per-sample quantum simulation cost: 4-15s per epoch on subsample.", 0),
], 0.5, 1.8, 12, 2.5, size=14)

add_text(slide, "Limitations", 0.5, 4.5, 12, 0.5, size=20, bold=True, color=MAROON)
add_bullets(slide, [
    ("Single train/test split (cross-validation prohibitively expensive).", 0),
    ("Noiseless simulator only; real quantum hardware would degrade.", 0),
    ("Small subgroups (AIAN N=36, Unknown N=18) limit fairness audit power.", 0),
    ("Right-censoring assumed non-informative.", 0),
    ("Cohort: SEER TNBC only, may not transfer to other cancers/subtypes.", 0),
], 0.5, 5, 12, 2.5, size=14)


# ============================================================
# Slide 14: Conclusion & Future Work
# ============================================================
slide = add_slide("Conclusion & Future Work")

add_text(slide, "Key takeaways", 0.5, 1.3, 12, 0.5, size=20, bold=True, color=TEAL)
add_bullets(slide, [
    ("Rigorous leakage auditing is essential: caught Year_of_follow_up_recode (AUC 0.98).", 0),
    ("Cox + Quantum Residual achieves C-index 0.7364, beating Cox PH 0.7326 (+0.0038).", 0),
    ("Mathematically guaranteed ≥ Cox PH performance via residual offset framework.", 0),
    ("Cross-framework consistency: residual @60mo AUC 0.7426 vs LightGBM 0.7462.", 0),
    ("Subgroup C-index gap 0.0754 (residual) vs 0.0842 (Cox PH alone).", 0),
], 0.5, 1.8, 12, 2.5, size=14)

add_text(slide, "Future work", 0.5, 4.5, 12, 0.5, size=20, bold=True, color=GOLD)
add_bullets(slide, [
    ("k-fold cross-validation when GPU-accelerated quantum simulators are available.", 0),
    ("Run on real quantum hardware (IBM Q, IonQ) to assess noise impact.", 0),
    ("Quantum kernel + Cox PH: drop variational training, use fixed feature maps.", 0),
    ("Extend to other cancer types and external validation cohorts.", 0),
    ("Fairness-constrained Cox loss for subgroup parity.", 0),
], 0.5, 5, 12, 2.5, size=14)


# ============================================================
# Slide 15: Thank you / References
# ============================================================
slide = add_slide("Thank you")

add_text(slide, "Questions?", 0.5, 2.5, 12.3, 1.0, size=44, bold=True,
         color=MAROON, align=PP_ALIGN.CENTER)

add_text(slide, "Code:  github.com/dataeducator/hybrid-classical-quantum",
         0.5, 4.0, 12.3, 0.5, size=18, color=BLUE, align=PP_ALIGN.CENTER)

add_text(slide, "Key references",
         0.5, 5.0, 12.3, 0.5, size=16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
refs = [
    "Cox, D.R. (1972). Regression Models and Life-Tables. JRSS-B 34:187-220.",
    "Katzman et al. (2018). DeepSurv. BMC Med Res Methodol 18:24.",
    "Pérez-Salinas et al. (2020). Data re-uploading universal classifier. Quantum 4:226.",
    "Qiu et al. (2024). TNBC survival prediction. Front Oncol 14:1388869.",
]
add_text(slide, refs, 0.5, 5.5, 12.3, 1.5, size=12, color=GRAY, align=PP_ALIGN.CENTER)


# Save
out_path = "Final_Presentation_Hybrid_QC_TNBC.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
